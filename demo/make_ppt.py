"""Generate hcompress demo PPT — minimalist light theme."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1a, 0x1a, 0x1a)
GRAY_800  = RGBColor(0x33, 0x33, 0x33)
GRAY_600  = RGBColor(0x66, 0x66, 0x66)
GRAY_400  = RGBColor(0x99, 0x99, 0x99)
GRAY_200  = RGBColor(0xe0, 0xe0, 0xe0)
GRAY_100  = RGBColor(0xf5, 0xf5, 0xf5)
BLUE      = RGBColor(0x4f, 0x6e, 0xf6)
GREEN     = RGBColor(0x1e, 0x9d, 0x5b)
RED       = RGBColor(0xd9, 0x40, 0x40)
TEAL      = RGBColor(0x1a, 0xab, 0x8e)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ──
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide

def add_text(slide, left, top, width, height, text, font_size=18,
             color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_line(slide, left, top, width, color=GRAY_200):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_placeholder(slide, left, top, width, height, label="插图位置"):
    """Dashed placeholder box for future images."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_100
    shape.line.color.rgb = GRAY_400
    shape.line.width = Pt(1)
    shape.line.dash_style = 2  # dash
    # Label text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY_400
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(6)
    return shape

def add_tag(slide, left, top, text, color=BLUE):
    """Small colored tag/badge."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.6), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

# ── Slide 1: Title ──
s = add_blank_slide()
add_line(s, 1.5, 3.5, 10.3, GRAY_200)
add_text(s, 1.5, 1.8, 10, 1.2, "hcompress", 72, BLACK, True)
add_text(s, 1.5, 2.7, 10, 0.6, "基于 Canonical Huffman 的高性能压缩软件", 24, GRAY_600)
add_text(s, 1.5, 3.8, 10, 0.5, "C 加速 · 插件架构 · 炸弹防护 · Electron 桌面端", 16, GRAY_400)
add_text(s, 1.5, 5.5, 10, 0.4, "Eric Huang  ·  2026.07", 13, GRAY_400)

# ── Slide 2: 问题 & 动机 ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 5, 0.6, "为什么做这个项目？", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
# Points
points = [
    ("📚", "算法落地", "哈夫曼编码是压缩领域的基础算法，但从理论到高性能实现中间隔着工程鸿沟"),
    ("🔌", "插件架构", "大部分压缩软件不支持扩展，无法按需添加炸弹检测、加密、格式支持"),
    ("🎨", "现代 UI", "开源压缩工具普遍界面老旧，缺乏暗色主题和拖拽交互"),
    ("⚡", "性能优化", "纯 Python 编解码速度慢，需要用 C 扩展加速热路径"),
]
for i, (icon, title, desc) in enumerate(points):
    y = 1.6 + i * 1.3
    add_text(s, 0.8, y, 0.5, 0.5, icon, 28)
    add_text(s, 1.4, y, 2, 0.4, title, 18, BLACK, True)
    add_text(s, 1.4, y + 0.35, 6, 0.6, desc, 13, GRAY_600)
add_placeholder(s, 8.2, 1.6, 4.5, 5.0, "📷 问题场景示意图")

# ── Slide 3: 功能特性 ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 5, 0.6, "核心功能", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
features = [
    ("Canonical Huffman", "业界标准算法，gzip/PNG 同款", BLUE),
    ("C 扩展加速", "编码 4.4× / 解码 29.5×", GREEN),
    ("10 个插件接口", "可替换熵编码、变换、校验、IO", BLUE),
    ("BombGuard", "零 IO 开销的炸弹检测 (100:1)", RED),
    ("多格式解压", "gzip / zip / tar / bz2 / xz", TEAL),
    ("文件夹归档", "目录 → 单文件 .hcf → 完整还原", GREEN),
    ("Electron GUI", "三主题切换 + 拖拽 + 插件面板", BLUE),
    ("69 测试", "全绿，覆盖边界 + 接口契约", GRAY_600),
]
for i, (title, desc, color) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.5 + row * 1.35
    add_tag(s, x, y, title, color)
    add_text(s, x + 1.8, y - 0.02, 3.8, 0.35, desc, 13, GRAY_600)
add_placeholder(s, 8.5, 0.8, 4.2, 6.2, "📷 功能截图拼贴")

# ── Slide 4: 架构总览 ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 8, 0.6, "系统架构", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
add_placeholder(s, 0.8, 1.6, 11.7, 5.4, "📷 架构图 — 建议用 demo/index.html 的 3D 架构截图")

# ── Slide 5: HCF 文件格式 ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 8, 0.6, "HCF 文件格式", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
fields = [
    ("Magic", "4B", "'HCF\\x1a'"),
    ("Version", "2B", "uint16 LE (v1 = 0x0001)"),
    ("Flags", "2B", "编码器ID / 压缩级别 / 扩展标志"),
    ("CRC-16", "2B", "Header 完整性校验"),
    ("Bit-Lengths", "256B", "每符号 1 字节，0 = 未出现"),
    ("Original Size", "8B", "原始文件大小 ← 炸弹检测锚点"),
    ("Ext Data", "可变", "JSON UTF-8，插件扩展数据"),
]
for i, (name, size, desc) in enumerate(fields):
    y = 1.5 + i * 0.65
    add_text(s, 0.8, y, 2.0, 0.4, name, 14, BLACK, True, font_name="Cascadia Code")
    add_text(s, 2.9, y, 1.0, 0.4, size, 13, GRAY_400, font_name="Cascadia Code")
    add_text(s, 4.0, y, 5.0, 0.4, desc, 13, GRAY_600)
add_text(s, 0.8, 6.0, 5, 0.3, "Header 固定 276 字节，紧凑且可扩展", 11, GRAY_400)
add_placeholder(s, 8.5, 1.6, 4.2, 5.2, "📷 Header 结构图")

# ── Slide 6: 性能数据 ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 10, 0.6, "性能实测", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
add_text(s, 0.8, 1.5, 5, 0.4, "C 扩展 vs 纯 Python — 530 KB 文本文件", 14, GRAY_600)
add_placeholder(s, 0.8, 2.1, 5.8, 2.8, "📷 编码/解码对比柱状图")
add_placeholder(s, 7.0, 2.1, 5.5, 2.8, "📷 不同文件类型压缩率对比")

perf_text = [
    "• C 扩展编码加速 4.4×（385ms → 88ms）",
    "• C 扩展解码加速 29.5×（809ms → 28ms）",
    "• 多进程压缩加速 3.1×（4 workers: 900ms → 291ms）",
    "• 文本压缩率稳定 50% 左右",
    "• 二进制/已压缩文件 — 几乎无效果（正常现象）",
    "• Python GIL 绕过：ProcessPoolExecutor + C 扩展",
]
for i, line in enumerate(perf_text):
    add_text(s, 0.8, 5.3 + i * 0.4, 11, 0.35, line, 13, GRAY_600)

# ── Slide 7: 多线程并行 ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 10, 0.6, "多进程并行压缩", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
add_placeholder(s, 0.8, 1.6, 6.0, 5.4, "📷 多线程架构图 / 性能对比图")
add_text(s, 7.5, 1.6, 5, 0.4, "ProcessPoolExecutor", 20, BLACK, True)
mt_points = [
    "• Python GIL 绕过：多进程替代多线程",
    "• 大文件自动分块，每块独立压缩",
    "• C 扩展 + 多进程 = 最大性能",
    "",
    "实测数据（5.2 MB 文本）：",
    "  单线程:  900 ms",
    "  2 进程:  377 ms (2.4×)",
    "  4 进程:  291 ms (3.1×)",
    "",
    "HCF 多块格式：",
    "  [Header] [BlockCount]",
    "  [Len₁][Data₁][Len₂][Data₂]...",
]
for i, txt in enumerate(mt_points):
    add_text(s, 7.5, 2.2 + i * 0.38, 5, 0.35, txt, 13, GRAY_800)

# ── Slide 8: 插件系统 ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 10, 0.6, "插件架构", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
add_placeholder(s, 0.8, 1.6, 6.5, 5.4, "📷 接口关系图 — 建议用 demo/index.html 的接口力导向图")
add_text(s, 7.8, 1.6, 5, 0.4, "10 个可扩展接口", 18, BLACK, True)
ifaces = [
    "IEntropyCodec — 熵编码算法",
    "ITransform — 数据变换 (BWT/MTF/RLE)",
    "IFilter — 预处理过滤器",
    "IMatchFinder — LZ77 字典匹配",
    "IChecksum — 完整性校验",
    "IIOBackend — IO 后端",
    "IBlockSplitter — 分块策略",
    "ICompressHook / IDecompressHook",
    "IObserver — 进度/日志",
    "★ IExtension — 万能自定义",
]
for i, txt in enumerate(ifaces):
    add_text(s, 7.8, 2.2 + i * 0.42, 5, 0.35, txt, 12, GRAY_800, font_name="Cascadia Code")
add_text(s, 7.8, 6.5, 5, 0.5, "插件 SDK: 4 行代码即可开发", 12, BLUE, True)

# ── Slide 8: BombGuard ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 10, 0.6, "BombGuard — 压缩包炸弹检测", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
add_placeholder(s, 7.5, 1.6, 5.0, 5.4, "📷 炸弹拦截弹窗截图")
add_text(s, 0.8, 1.6, 6, 0.4, "原理", 20, BLACK, True)
steps = [
    "1. HCF Header 中存储 original_size（原始文件大小）",
    "2. 解压前读取 header → 仅 276 字节，零 IO 开销",
    "3. 计算膨胀比 = original_size / 压缩文件大小",
    "4. 若 > 100:1 → 判断为炸弹 → 拒绝解压",
    "5. 递归深度检测 → 防嵌套炸弹（.hcf 套 .hcf）",
]
for i, txt in enumerate(steps):
    add_text(s, 0.8, 2.2 + i * 0.55, 6, 0.4, txt, 15, GRAY_800)
add_text(s, 0.8, 5.2, 6, 0.5, "实际测试：291 字节文件声称 10GB → 膨胀比 3689 万:1 → 被拦截", 13, RED, True)

# ── Slide 9: UI 展示 ──
s = add_blank_slide()
add_text(s, 0.8, 0.5, 10, 0.6, "v2 桌面端 — Electron + React", 32, BLACK, True)
add_line(s, 0.8, 1.2, 11.7, GRAY_200)
add_placeholder(s, 0.8, 1.6, 7.8, 5.4, "📷 v2 GUI 完整截图（含三主题对比）")
features_v2 = [
    "☀️/💻/🌙 三主题切换",
    "拖拽上传 + 文件列表",
    "压缩率实时估算",
    "插件管理面板（开关 + 状态）",
    "BombGuard 可视化拦截弹窗",
    "归档浏览器",
    "F11 全屏模式",
]
for i, txt in enumerate(features_v2):
    add_text(s, 9.0, 1.8 + i * 0.55, 3.8, 0.4, txt, 14, GRAY_800)

# ── Slide 10: 项目信息 & 感谢 ──
s = add_blank_slide()
add_line(s, 1.5, 3.8, 10.3, GRAY_200)
add_text(s, 1.5, 1.5, 10, 1.2, "Thank You", 72, BLACK, True)
add_text(s, 1.5, 2.5, 10, 0.6, "hcompress — Canonical Huffman 压缩工具", 22, GRAY_600)
info = [
    "GitHub v1: github.com/Eric-huang799/hcompress",
    "GitHub v2: github.com/Eric-huang799/hcompress-v2",
    "技术栈: Python + C + Electron + React + TypeScript",
    "测试: 69 tests · 性能: 编码 4.4× 解码 29.5×",
]
for i, txt in enumerate(info):
    add_text(s, 1.5, 4.2 + i * 0.45, 10, 0.4, txt, 15, GRAY_600)
add_text(s, 1.5, 6.5, 10, 0.4, "Eric Huang · 2026.07", 13, GRAY_400)

# ── Save ──
out = os.path.join(os.path.dirname(__file__), "hcompress-demo.pptx")
prs.save(out)
print(f"PPT saved: {out}")
print(f"Slides: {len(prs.slides)}")
